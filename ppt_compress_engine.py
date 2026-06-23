"""PPT 压缩引擎 —— 纯逻辑，无界面依赖。

核心思路：.pptx 本质是 zip，体积几乎都在 ppt/media/ 里。
本引擎把图片按"在幻灯片上的实际显示尺寸 × 目标 DPI"降到合理分辨率，
再按档位重新编码 JPEG。视频/音频/嵌入字体 v1 一律不动（最高损伤/风险）。

最小损伤优先级（档位越靠前越温和）：
  1. 只缩"超规格"的大图（显示很小却存了超高分辨率）——低损
  2. 适度降低 JPEG 质量 ——中损
  3. 视频、嵌入字体 —— v1 不碰

设计要点：
- 重压后的图片**保持原文件名与原格式**（png 仍存 png，jpg 仍存 jpg），
  这样幻灯片里的引用关系、内容类型表都不用改，最稳。
- 只在"压完更小"时才替换，绝不把文件改大。
- 整个新 zip 在内存里构建，估算(预览)和实际保存共用同一套代码。
"""

from __future__ import annotations

import hashlib
import io
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image

EMU_PER_INCH = 914400

# ============ 质量档位（滑块就映射到这里） ============
# (显示名, 目标DPI, JPEG质量, 一句话说明)
QUALITY_TIERS = [
    ("几乎无损", 220, 92, "只缩极端超规格的大图，质量几乎不变，最大限度保留观看体验"),
    ("高质量", 160, 88, "适合投影/演示，肉眼基本无差别"),
    ("标准", 120, 82, "体积与清晰度的平衡点（推荐）"),
    ("压缩优先", 96, 74, "明显减小体积，近看略有损失"),
    ("最小体积", 72, 62, "最大压缩，适合网络传输/存档"),
]
DEFAULT_TIER = 2

# 只重压这些位图格式；矢量(emf/wmf)、视频、音频等保持原样
_RECOMPRESS_EXT = {".jpg", ".jpeg", ".png"}
_IMAGE_EXT = _RECOMPRESS_EXT | {".gif", ".bmp", ".tiff", ".tif"}
_VIDEO_EXT = {".mp4", ".mov", ".m4v", ".avi", ".wmv", ".mkv", ".webm"}
_AUDIO_EXT = {".mp3", ".wav", ".m4a", ".aac", ".wma"}


@dataclass
class MediaInfo:
    name: str            # zip 内路径，如 ppt/media/image1.png
    size: int
    kind: str            # image | video | audio | font | other


@dataclass
class AnalyzeResult:
    path: Path
    total_size: int
    media_size: int
    image_size: int
    video_size: int
    audio_size: int
    image_count: int
    video_count: int
    has_embedded_fonts: bool
    media: list[MediaInfo] = field(default_factory=list)

    @property
    def compressible_size(self) -> int:
        """v1 实际能动的部分（图片）。"""
        return self.image_size


@dataclass
class CompressResult:
    original_size: int
    new_size: int
    images_processed: int
    images_shrunk: int
    skipped_video_size: int     # 未被压缩的视频总字节（关掉视频压缩或无 ffmpeg 时）
    videos_processed: int = 0
    videos_shrunk: int = 0
    data: bytes | None = None   # 估算时为 None，保存时回填以便写盘

    @property
    def saved(self) -> int:
        return max(0, self.original_size - self.new_size)

    @property
    def ratio(self) -> float:
        if self.original_size <= 0:
            return 0.0
        return self.saved / self.original_size


def _kind_for(name: str) -> str:
    n = name.lower()
    ext = Path(n).suffix
    if n.startswith("ppt/fonts/"):
        return "font"
    if ext in _IMAGE_EXT:
        return "image"
    if ext in _VIDEO_EXT:
        return "video"
    if ext in _AUDIO_EXT:
        return "audio"
    return "other"


def analyze(path: str | Path) -> AnalyzeResult:
    """快速扫描 pptx，统计各类媒体体积，用于界面展示与提示。"""
    path = Path(path)
    total = path.stat().st_size
    media: list[MediaInfo] = []
    image_size = video_size = audio_size = media_total = 0
    image_count = video_count = 0
    has_fonts = False

    with zipfile.ZipFile(path) as zf:
        for info in zf.infolist():
            name = info.filename
            kind = _kind_for(name)
            if name.startswith("ppt/fonts/"):
                has_fonts = True
            if not (name.startswith("ppt/media/") or kind == "font"):
                continue
            size = info.file_size
            media_total += size
            media.append(MediaInfo(name=name, size=size, kind=kind))
            if kind == "image":
                image_size += size
                image_count += 1
            elif kind == "video":
                video_size += size
                video_count += 1
            elif kind == "audio":
                audio_size += size

    return AnalyzeResult(
        path=path,
        total_size=total,
        media_size=media_total,
        image_size=image_size,
        video_size=video_size,
        audio_size=audio_size,
        image_count=image_count,
        video_count=video_count,
        has_embedded_fonts=has_fonts,
        media=media,
    )


def _build_display_map(path: Path) -> dict[str, int]:
    """用 python-pptx 遍历所有图片形状，得到 sha1 -> 最大显示长边(EMU)。

    一张图可能在多处以不同大小使用，取最大显示尺寸，避免把"某处放得很大"的图压糊。
    遍历范围覆盖幻灯片、版式、母版。拿不到 python-pptx 时返回空表（退化为按整页估算）。
    """
    try:
        from pptx import Presentation
    except Exception:
        return {}

    display: dict[str, int] = {}

    def visit(shapes):
        for shape in shapes:
            # 组合形状：递归
            if getattr(shape, "shape_type", None) == 6:  # GROUP
                try:
                    visit(shape.shapes)
                except Exception:
                    pass
                continue
            try:
                image = shape.image  # 仅图片形状可取，其它抛异常
            except Exception:
                continue
            try:
                long_edge = max(int(shape.width or 0), int(shape.height or 0))
            except Exception:
                long_edge = 0
            if long_edge <= 0:
                continue
            sha1 = image.sha1
            if long_edge > display.get(sha1, 0):
                display[sha1] = long_edge

    try:
        prs = Presentation(str(path))
    except Exception:
        return {}

    for slide in prs.slides:
        visit(slide.shapes)
    for master in prs.slide_masters:
        visit(master.shapes)
        for layout in master.slide_layouts:
            visit(layout.shapes)

    return display


def _slide_long_edge_emu(path: Path) -> int:
    """整页长边(EMU)，作为显示尺寸未知时的保守上限。"""
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        return max(int(prs.slide_width or 0), int(prs.slide_height or 0)) or 12192000
    except Exception:
        return 12192000  # 16:9 默认 13.33in


def _recompress_image(
    data: bytes,
    display_emu: int | None,
    slide_long_emu: int,
    target_dpi: int,
    jpeg_quality: int,
) -> bytes | None:
    """重压单张图片，返回新字节；无法处理或没变小则返回 None。"""
    try:
        im = Image.open(io.BytesIO(data))
        im.load()
    except Exception:
        return None

    fmt = (im.format or "").upper()
    if fmt not in ("JPEG", "PNG"):
        return None

    # 目标像素：显示英寸 × DPI；显示尺寸未知时按整页长边（保守，不会过度缩小）
    disp_emu = display_emu if display_emu else slide_long_emu
    target_px = max(1, round(disp_emu / EMU_PER_INCH * target_dpi))

    w, h = im.size
    long_edge = max(w, h)
    resized = False
    if long_edge > target_px:
        scale = target_px / long_edge
        im = im.resize((max(1, round(w * scale)), max(1, round(h * scale))), Image.LANCZOS)
        resized = True

    out = io.BytesIO()
    if fmt == "JPEG":
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        im.save(out, "JPEG", quality=jpeg_quality, optimize=True, progressive=True)
    else:  # PNG —— 保持格式与透明通道，靠降分辨率 + optimize 减小
        if not resized:
            # 没缩放就只能靠 optimize，收益有限；仍尝试，由调用方决定是否采用
            im.save(out, "PNG", optimize=True)
        else:
            im.save(out, "PNG", optimize=True)

    return out.getvalue()


def compress_to_bytes(
    path: str | Path,
    tier_index: int = DEFAULT_TIER,
    keep_data: bool = True,
    progress_cb=None,
    compress_video: bool = False,
    ffmpeg: str | None = None,
    video_estimate: bool = False,
) -> CompressResult:
    """把整个 pptx 重压到内存，返回结果（含新字节，便于保存或测量）。

    progress_cb(done, total) 可选，用于界面进度。
    compress_video: 是否同时压缩视频（需要 ffmpeg 路径）。
    video_estimate: 预览模式——视频不真正转码，按码率快速估算大小（更快）。
    """
    path = Path(path)
    _, target_dpi, jpeg_quality, _ = QUALITY_TIERS[tier_index]

    do_video = compress_video and ffmpeg
    if do_video:
        import video as videomod

    display_map = _build_display_map(path)
    slide_long = _slide_long_edge_emu(path)

    images_processed = images_shrunk = 0
    videos_processed = videos_shrunk = 0
    skipped_video = 0
    original_size = path.stat().st_size

    src = zipfile.ZipFile(path)
    names = src.namelist()
    total = len(names)

    out_buf = io.BytesIO()
    with zipfile.ZipFile(out_buf, "w", zipfile.ZIP_DEFLATED) as dst:
        for i, name in enumerate(names):
            data = src.read(name)
            lname = name.lower()
            ext = Path(lname).suffix

            if lname.startswith("ppt/media/"):
                if ext in _RECOMPRESS_EXT:
                    images_processed += 1
                    sha1 = hashlib.sha1(data).hexdigest()
                    display_emu = display_map.get(sha1)
                    new = _recompress_image(
                        data, display_emu, slide_long, target_dpi, jpeg_quality
                    )
                    if new is not None and len(new) < len(data):
                        data = new
                        images_shrunk += 1
                elif ext in _VIDEO_EXT:
                    if do_video and ext in videomod.SUPPORTED_VIDEO_EXT:
                        videos_processed += 1
                        if video_estimate:
                            # 预览：不真正写入转码结果，只调整记账用的大小
                            est = videomod.estimate_size(data, ext, ffmpeg, tier_index)
                            if est < len(data):
                                videos_shrunk += 1
                                # 写入占位的同尺寸数据仅为让 new_size 反映估算值
                                data = data[:est] if est <= len(data) else data
                        else:
                            new = videomod.transcode(data, ext, ffmpeg, tier_index)
                            if new is not None:
                                data = new
                                videos_shrunk += 1
                            else:
                                skipped_video += len(data)
                    else:
                        skipped_video += len(data)

            dst.writestr(name, data)
            if progress_cb:
                progress_cb(i + 1, total)

    src.close()

    new_bytes = out_buf.getvalue()
    return CompressResult(
        original_size=original_size,
        new_size=len(new_bytes),
        images_processed=images_processed,
        images_shrunk=images_shrunk,
        skipped_video_size=skipped_video,
        videos_processed=videos_processed,
        videos_shrunk=videos_shrunk,
        data=new_bytes if keep_data else None,
    )


def estimate(
    path: str | Path,
    tier_index: int = DEFAULT_TIER,
    progress_cb=None,
    compress_video: bool = False,
    ffmpeg: str | None = None,
) -> CompressResult:
    """只测量压缩后大小，不保留数据（用于滑块松手时的真实预览）。

    视频用码率快速估算，不真正转码，避免界面卡死。
    """
    return compress_to_bytes(
        path, tier_index, keep_data=False, progress_cb=progress_cb,
        compress_video=compress_video, ffmpeg=ffmpeg, video_estimate=True,
    )


def compress_file(
    path: str | Path,
    out_path: str | Path,
    tier_index: int = DEFAULT_TIER,
    progress_cb=None,
    compress_video: bool = False,
    ffmpeg: str | None = None,
) -> CompressResult:
    """压缩并写入 out_path。视频若开启则真正转码。"""
    result = compress_to_bytes(
        path, tier_index, keep_data=True, progress_cb=progress_cb,
        compress_video=compress_video, ffmpeg=ffmpeg, video_estimate=False,
    )
    Path(out_path).write_bytes(result.data)
    result.data = None  # 写完释放
    return result


def human_size(num: int) -> str:
    value = float(num)
    for unit in ("B", "KB", "MB", "GB"):
        if value < 1024 or unit == "GB":
            return f"{value:.0f} {unit}" if unit == "B" else f"{value:.1f} {unit}"
        value /= 1024
    return f"{value:.1f} GB"
